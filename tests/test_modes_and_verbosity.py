from __future__ import annotations

import json
import os
import subprocess
import sys

from pptx.util import Inches

from report_comparator import compare_runs

from .fixture_builder import DeckBuilder, make_png
from .test_walking_skeleton import save_pair


def test_lenient_mode_accepts_minor_findings_without_status_escalation(tmp_path):
    image_path = make_png(tmp_path / "pic.png")
    old = DeckBuilder()
    new = DeckBuilder()
    old.add_picture("WaferMap", image_path, left=1, top=1)
    new.add_picture("WaferMap", image_path, left=1.2, top=1)
    old_dir, new_dir = save_pair(tmp_path, old, new)

    strict = compare_runs(old_dir, new_dir, {"shift_tolerance": int(Inches(0.1))})
    lenient = compare_runs(old_dir, new_dir, {"mode": "lenient", "shift_tolerance": int(Inches(0.1))})

    assert strict["summary"] == {"files": 1, "ok": 0, "warnings": 1, "failures": 0}
    assert strict["files"][0]["slides"][0]["findings"][0]["severity"] == "minor"
    assert lenient["summary"] == {"files": 1, "ok": 1, "warnings": 0, "failures": 0}
    assert lenient["files"][0]["slides"][0]["status"] == "ok"
    assert lenient["files"][0]["slides"][0]["findings"][0]["severity"] == "accepted"


def test_quiet_suppresses_minor_and_accepted_entries_but_keeps_failures(tmp_path):
    image_path = make_png(tmp_path / "pic.png")
    old = DeckBuilder()
    new = DeckBuilder()
    old.add_picture("WaferMap", image_path, left=1, top=1)
    new.add_picture("WaferMap", image_path, left=1.2, top=1)
    old_dir, new_dir = save_pair(tmp_path, old, new)

    quiet_minor = compare_runs(old_dir, new_dir, {"quiet": True, "shift_tolerance": int(Inches(0.1))})

    assert quiet_minor["summary"] == {"files": 1, "ok": 1, "warnings": 0, "failures": 0}
    assert quiet_minor["files"][0]["slides"] == []

    failing_old = DeckBuilder()
    failing_new = DeckBuilder()
    failing_old.add_text("Description", "Wafer map centered")
    failing_new.add_text("Description", "Wafer map shifted")
    failure_dir = tmp_path / "failure"
    failure_dir.mkdir()
    old_dir, new_dir = save_pair(failure_dir, failing_old, failing_new)

    quiet_fail = compare_runs(old_dir, new_dir, {"quiet": True})

    assert quiet_fail["summary"] == {"files": 1, "ok": 0, "warnings": 0, "failures": 1}
    assert quiet_fail["files"][0]["slides"][0]["findings"][0]["severity"] == "fail"


def test_calibrate_emits_measured_metrics_for_ok_elements(tmp_path):
    image_path = make_png(tmp_path / "pic.png")
    old = DeckBuilder()
    new = DeckBuilder()
    old.add_picture("WaferMap", image_path, left=1, top=1)
    new.add_picture("WaferMap", image_path, left=1, top=1)
    old_dir, new_dir = save_pair(tmp_path, old, new)

    normal = compare_runs(old_dir, new_dir, {})
    calibrated = compare_runs(old_dir, new_dir, {"calibrate": True})

    assert normal["files"][0]["slides"] == []
    assert calibrated["summary"] == {"files": 1, "ok": 1, "warnings": 0, "failures": 0}
    metric = calibrated["files"][0]["slides"][0]["findings"][0]
    assert metric["severity"] == "ok"
    assert metric["type"] == "measured_metrics"
    assert metric["element"] == "picture:WaferMap"
    assert metric["metrics"]["shift_delta"] == 0
    assert metric["metrics"]["picture_changed_percent"] == 0.0


def test_cli_mode_override_and_quiet_flag_control_minor_output(tmp_path):
    image_path = make_png(tmp_path / "pic.png")
    old = DeckBuilder()
    new = DeckBuilder()
    old.add_picture("WaferMap", image_path, left=1, top=1)
    new.add_picture("WaferMap", image_path, left=1.2, top=1)
    old_dir, new_dir = save_pair(tmp_path, old, new)
    config = tmp_path / "config.yaml"
    config.write_text(f"mode: lenient\nshift_tolerance: {int(Inches(0.1))}\n", encoding="utf-8")

    strict_out = tmp_path / "strict.json"
    subprocess.run(
        [
            sys.executable,
            "-m",
            "report_comparator.cli",
            "--old",
            str(old_dir),
            "--new",
            str(new_dir),
            "--config",
            str(config),
            "--mode",
            "strict",
            "--out",
            str(strict_out),
        ],
        check=True,
        env={**os.environ, "PYTHONPATH": "src"},
    )
    strict_report = json.loads(strict_out.read_text(encoding="utf-8"))
    assert strict_report["summary"] == {"files": 1, "ok": 0, "warnings": 1, "failures": 0}
    assert strict_report["files"][0]["slides"][0]["findings"][0]["severity"] == "minor"

    quiet_out = tmp_path / "quiet.json"
    subprocess.run(
        [
            sys.executable,
            "-m",
            "report_comparator.cli",
            "--old",
            str(old_dir),
            "--new",
            str(new_dir),
            "--config",
            str(config),
            "--quiet",
            "--out",
            str(quiet_out),
        ],
        check=True,
        env={**os.environ, "PYTHONPATH": "src"},
    )
    quiet_report = json.loads(quiet_out.read_text(encoding="utf-8"))
    assert quiet_report["summary"] == {"files": 1, "ok": 1, "warnings": 0, "failures": 0}
    assert quiet_report["files"][0]["slides"] == []
