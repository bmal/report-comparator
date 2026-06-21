from __future__ import annotations

from pptx.dml.color import RGBColor

from report_comparator import compare_runs

from .fixture_builder import DeckBuilder
from .test_walking_skeleton import save_pair


def test_recolored_text_with_same_words_is_minor(tmp_path):
    old = DeckBuilder()
    new = DeckBuilder()
    old.add_text("Description", "Wafer map centered")
    text = new.add_text("Description", "Wafer map centered")
    text.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 0, 0)
    old_dir, new_dir = save_pair(tmp_path, old, new)

    report = compare_runs(old_dir, new_dir, {})

    finding = report["files"][0]["slides"][0]["findings"][0]
    assert report["summary"] == {"files": 1, "ok": 0, "warnings": 1, "failures": 0}
    assert finding["severity"] == "minor"
    assert finding["type"] == "text_formatting_changed"


def test_bolded_table_cell_with_same_value_is_minor(tmp_path):
    old = DeckBuilder()
    new = DeckBuilder()
    old.add_table("Data", values=[["1.23"]])
    table = new.add_table("Data", values=[["1.23"]])
    table.table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.bold = True
    old_dir, new_dir = save_pair(tmp_path, old, new)

    report = compare_runs(old_dir, new_dir, {})

    finding = report["files"][0]["slides"][0]["findings"][0]
    assert report["summary"] == {"files": 1, "ok": 0, "warnings": 1, "failures": 0}
    assert finding["severity"] == "minor"
    assert finding["type"] == "table_formatting_changed"


def test_changed_table_cell_fill_with_same_value_is_minor(tmp_path):
    old = DeckBuilder()
    new = DeckBuilder()
    old.add_table("Data", values=[["1.23"]])
    table = new.add_table("Data", values=[["1.23"]])
    cell = table.table.cell(0, 0)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(255, 255, 0)
    old_dir, new_dir = save_pair(tmp_path, old, new)

    report = compare_runs(old_dir, new_dir, {})

    finding = report["files"][0]["slides"][0]["findings"][0]
    assert report["summary"] == {"files": 1, "ok": 0, "warnings": 1, "failures": 0}
    assert finding["severity"] == "minor"
    assert finding["type"] == "table_formatting_changed"


def test_formatting_drift_does_not_suppress_content_failure(tmp_path):
    old = DeckBuilder()
    new = DeckBuilder()
    old.add_text("Description", "Wafer map centered")
    text = new.add_text("Description", "Wafer map shifted")
    text.text_frame.paragraphs[0].runs[0].font.bold = True
    old_dir, new_dir = save_pair(tmp_path, old, new)

    report = compare_runs(old_dir, new_dir, {})

    findings = report["files"][0]["slides"][0]["findings"]
    assert report["summary"] == {"files": 1, "ok": 0, "warnings": 0, "failures": 1}
    assert {finding["type"] for finding in findings} == {"text_changed"}
    assert findings[0]["severity"] == "fail"
