from __future__ import annotations

from pptx.util import Inches

from report_comparator import compare_runs

from .fixture_builder import DeckBuilder, make_png
from .test_walking_skeleton import save_pair


def test_shift_tolerance_controls_ok_vs_minor_placement_drift(tmp_path):
    image_path = make_png(tmp_path / "pic.png")
    old = DeckBuilder()
    new = DeckBuilder()
    old.add_picture("WaferMap", image_path, left=1, top=1)
    new.add_picture("WaferMap", image_path, left=1.2, top=1)
    old_dir, new_dir = save_pair(tmp_path, old, new)

    tolerated = compare_runs(old_dir, new_dir, {"shift_tolerance": int(Inches(0.25))})
    flagged = compare_runs(old_dir, new_dir, {"shift_tolerance": int(Inches(0.1))})

    assert tolerated["summary"] == {"files": 1, "ok": 1, "warnings": 0, "failures": 0}
    finding = flagged["files"][0]["slides"][0]["findings"][0]
    assert flagged["summary"] == {"files": 1, "ok": 0, "warnings": 1, "failures": 0}
    assert finding["severity"] == "minor"
    assert finding["type"] == "placement_changed"


def test_resize_tolerance_controls_ok_vs_minor_placement_drift(tmp_path):
    image_path = make_png(tmp_path / "pic.png")
    old = DeckBuilder()
    new = DeckBuilder()
    old.add_picture("WaferMap", image_path, left=1, top=1, width=1)
    new.add_picture("WaferMap", image_path, left=1, top=1, width=1.2)
    old_dir, new_dir = save_pair(tmp_path, old, new)

    tolerated = compare_runs(old_dir, new_dir, {"resize_tolerance": int(Inches(0.25))})
    flagged = compare_runs(old_dir, new_dir, {"resize_tolerance": int(Inches(0.1))})

    assert tolerated["summary"] == {"files": 1, "ok": 1, "warnings": 0, "failures": 0}
    finding = flagged["files"][0]["slides"][0]["findings"][0]
    assert flagged["summary"] == {"files": 1, "ok": 0, "warnings": 1, "failures": 0}
    assert finding["severity"] == "minor"
    assert finding["type"] == "placement_changed"


def test_element_extending_past_slide_bounds_is_failure(tmp_path):
    image_path = make_png(tmp_path / "pic.png")
    old = DeckBuilder()
    new = DeckBuilder()
    old.add_picture("WaferMap", image_path, left=1, top=1)
    new.add_picture("WaferMap", image_path, left=9.5, top=1, width=1)
    old_dir, new_dir = save_pair(tmp_path, old, new)

    report = compare_runs(old_dir, new_dir, {"shift_tolerance": int(Inches(20))})

    finding = report["files"][0]["slides"][0]["findings"][0]
    assert report["summary"] == {"files": 1, "ok": 0, "warnings": 0, "failures": 1}
    assert finding["severity"] == "fail"
    assert finding["type"] == "off_slide"


def test_new_non_picture_overlap_fails_but_existing_reference_overlap_is_allowed(tmp_path):
    old = DeckBuilder()
    new = DeckBuilder()
    old.add_text("LabelA", "same", left=1, top=1, width=1)
    old.add_text("LabelB", "same", left=3, top=1, width=1)
    new.add_text("LabelA", "same", left=1, top=1, width=1)
    new.add_text("LabelB", "same", left=1.5, top=1, width=1)
    old_dir, new_dir = save_pair(tmp_path, old, new)

    new_overlap = compare_runs(old_dir, new_dir, {"shift_tolerance": int(Inches(20))})

    findings = new_overlap["files"][0]["slides"][0]["findings"]
    assert new_overlap["summary"] == {"files": 1, "ok": 0, "warnings": 0, "failures": 1}
    assert any(finding["severity"] == "fail" and finding["type"] == "new_overlap" for finding in findings)

    old_overlapping = DeckBuilder()
    new_overlapping = DeckBuilder()
    old_overlapping.add_text("LabelA", "same", left=1, top=1, width=1)
    old_overlapping.add_text("LabelB", "same", left=1.5, top=1, width=1)
    new_overlapping.add_text("LabelA", "same", left=1, top=1, width=1)
    new_overlapping.add_text("LabelB", "same", left=1.5, top=1, width=1)
    existing_dir = tmp_path / "existing"
    existing_dir.mkdir()
    old_dir, new_dir = save_pair(existing_dir, old_overlapping, new_overlapping)

    existing_overlap = compare_runs(old_dir, new_dir, {})

    assert existing_overlap["summary"] == {"files": 1, "ok": 1, "warnings": 0, "failures": 0}
    assert existing_overlap["files"][0]["slides"] == []


def test_picture_to_picture_overlap_fails_even_when_reference_already_overlapped(tmp_path):
    image_path = make_png(tmp_path / "pic.png")
    old = DeckBuilder()
    new = DeckBuilder()
    old.add_picture("MapA", image_path, left=1, top=1, width=1)
    old.add_picture("MapB", image_path, left=1.5, top=1, width=1)
    new.add_picture("MapA", image_path, left=1, top=1, width=1)
    new.add_picture("MapB", image_path, left=1.5, top=1, width=1)
    old_dir, new_dir = save_pair(tmp_path, old, new)

    report = compare_runs(old_dir, new_dir, {})

    finding = report["files"][0]["slides"][0]["findings"][0]
    assert report["summary"] == {"files": 1, "ok": 0, "warnings": 0, "failures": 1}
    assert finding["severity"] == "fail"
    assert finding["type"] == "picture_overlap"
