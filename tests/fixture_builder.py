from __future__ import annotations

from pathlib import Path

from PIL import Image, PngImagePlugin
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches


class DeckBuilder:
    def __init__(self) -> None:
        self.presentation = Presentation()
        self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

    @property
    def slide(self):
        return self.presentation.slides[0]

    def add_text(
        self, name: str, text: str = "Hello", left: float = 1, top: float = 1, width: float = 2, height: float = 0.5
    ):
        box = self.slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        box.name = name
        box.text = text
        return box

    def add_table(
        self,
        name: str,
        text: str = "value",
        left: float = 1,
        top: float = 2,
        width: float = 2,
        height: float = 0.5,
        values: list[list[str]] | None = None,
    ):
        values = values or [[text]]
        table_shape = self.slide.shapes.add_table(len(values), len(values[0]), Inches(left), Inches(top), Inches(width), Inches(height))
        table_shape.name = name
        for row_index, row in enumerate(values):
            for column_index, value in enumerate(row):
                table_shape.table.cell(row_index, column_index).text = value
        return table_shape

    def add_picture(
        self,
        name: str | None,
        image_path: Path,
        left: float = 1,
        top: float = 3,
        width: float = 1,
        height: float = 1,
        alt_text: str | None = None,
    ):
        picture = self.slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), Inches(width), Inches(height))
        if name is not None:
            picture.name = name
        if alt_text is not None:
            set_alt_text(picture, alt_text)
        return picture

    def add_group(
        self,
        text_name: str,
        text: str,
        picture_name: str,
        image_path: Path,
        text_pos: tuple[float, float, float, float] = (1, 1, 3, 0.5),
        picture_pos: tuple[float, float, float, float] = (1, 2, 1, 1),
    ):
        group = self.slide.shapes.add_group_shape()
        box = group.shapes.add_textbox(*(Inches(value) for value in text_pos))
        box.name = text_name
        box.text = text
        picture = group.shapes.add_picture(str(image_path), *(Inches(value) for value in picture_pos))
        picture.name = picture_name
        return group

    def add_unknown(self, name: str, left: float = 4, top: float = 1, width: float = 1, height: float = 1) -> None:
        shape = self.slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.name = name

    def add_slide(self) -> None:
        self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

    def save(self, path: Path) -> None:
        self.presentation.save(path)


def set_alt_text(shape, descr: str) -> None:
    c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
    c_nv_pr.set("descr", descr)


def make_png(
    path: Path,
    color: tuple[int, int, int] = (255, 0, 0),
    size: tuple[int, int] = (10, 10),
    changed_pixels: dict[tuple[int, int], tuple[int, int, int]] | None = None,
    metadata: dict[str, str] | None = None,
) -> Path:
    image = Image.new("RGB", size, color)
    for xy, pixel in (changed_pixels or {}).items():
        image.putpixel(xy, pixel)
    pnginfo = None
    if metadata:
        pnginfo = PngImagePlugin.PngInfo()
        for key, value in metadata.items():
            pnginfo.add_text(key, value)
    image.save(path, pnginfo=pnginfo)
    return path
