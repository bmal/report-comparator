from __future__ import annotations

from collections import defaultdict
from io import BytesIO
from pathlib import Path
import re
from typing import Any, Iterable

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def compare_runs(old_dir: str | Path, new_dir: str | Path, config: dict[str, Any] | None = None) -> dict[str, Any]:
    config = config or {}
    old_dir = Path(old_dir)
    new_dir = Path(new_dir)
    files: list[dict[str, Any]] = []

    old_pairs, old_errors = _collect_pairs(old_dir, config)
    new_pairs, new_errors = _collect_pairs(new_dir, config)
    files.extend(old_errors)
    files.extend(new_errors)

    old_keys = set(old_pairs)
    new_keys = set(new_pairs)
    for key in sorted(old_keys - new_keys):
        files.append(_file_error(key, str(old_pairs[key]), None, f"unmatched old deck for key '{key}'"))
    for key in sorted(new_keys - old_keys):
        files.append(_file_error(key, None, str(new_pairs[key]), f"unmatched new deck for key '{key}'"))

    for key in sorted(old_keys & new_keys):
        files.append(_compare_decks(key, old_pairs[key], new_pairs[key], config))

    return _with_summary(files)


def _collect_pairs(directory: Path, config: dict[str, Any]) -> tuple[dict[str, Path], list[dict[str, Any]]]:
    pattern = re.compile(config.get("filename_key_pattern", r"^(?P<key>.+)_\d{8}(_\d{6})?\.pptx$"))
    grouped: dict[str, list[Path]] = defaultdict(list)
    errors: list[dict[str, Any]] = []

    for path in sorted(directory.glob("*.pptx")):
        match = pattern.match(path.name)
        key = match.group("key") if match else path.stem
        grouped[key].append(path)

    pairs: dict[str, Path] = {}
    for key, paths in grouped.items():
        if len(paths) == 1:
            pairs[key] = paths[0]
        else:
            errors.append(_file_error(key, None, None, f"non-unique deck key '{key}' in {directory}"))
    return pairs, errors


def _compare_decks(key: str, old_path: Path, new_path: Path, config: dict[str, Any]) -> dict[str, Any]:
    file_result = {
        "key": key,
        "old": str(old_path),
        "new": str(new_path),
        "status": "ok",
        "errors": [],
        "slides": [],
    }
    try:
        old_deck = Presentation(str(old_path))
        new_deck = Presentation(str(new_path))
    except Exception as exc:  # python-pptx exposes several package/read errors.
        file_result["status"] = "fail"
        file_result["errors"].append({"severity": "fail", "type": "file_error", "message": f"could not read deck: {exc}"})
        return file_result

    old_count = len(old_deck.slides)
    new_count = len(new_deck.slides)
    if old_count != new_count:
        file_result["status"] = "fail"
        file_result["errors"].append(
            {
                "severity": "fail",
                "type": "slide_count_mismatch",
                "message": f"slide count changed from {old_count} to {new_count}",
            }
        )

    for index in range(min(old_count, new_count)):
        slide = _compare_slides(
            index + 1,
            old_deck.slides[index],
            new_deck.slides[index],
            int(new_deck.slide_width),
            int(new_deck.slide_height),
            config,
        )
        if slide["findings"]:
            file_result["slides"].append(slide)

    if any(f["severity"] == "fail" for slide in file_result["slides"] for f in slide["findings"]):
        file_result["status"] = "fail"
    elif any(f["severity"] == "minor" for slide in file_result["slides"] for f in slide["findings"]):
        file_result["status"] = "warning"
    return file_result


def _compare_slides(
    index: int, old_slide: Any, new_slide: Any, slide_width: int, slide_height: int, config: dict[str, Any]
) -> dict[str, Any]:
    old_elements = _visible_elements(old_slide.shapes, config)
    new_elements = _visible_elements(new_slide.shapes, config)
    findings: list[dict[str, Any]] = []
    matched, missing, stray = _match_elements(old_elements, new_elements)

    for element in missing:
        findings.append(_finding("fail", "missing_element", element["element"], f"missing {element['kind']} '{element['name']}'"))
    for element in stray:
        findings.append(_finding("fail", "stray_element", element["element"], f"stray {element['kind']} '{element['name']}'"))

    for old_element, new_element in matched:
        if config.get("calibrate", False):
            findings.append(_metric_finding(old_element, new_element))
        placement_finding = _compare_placement(old_element, new_element, slide_width, slide_height, config)
        if placement_finding:
            findings.append(placement_finding)
        if old_element["kind"] == "picture":
            finding = _compare_picture_content(old_element, new_element, config)
            if finding:
                findings.append(finding)
        if old_element["kind"] == "text":
            finding = _compare_text_content(old_element, new_element, config)
            if finding:
                findings.append(finding)
            else:
                formatting_finding = _compare_text_formatting(old_element, new_element)
                if formatting_finding:
                    findings.append(formatting_finding)
        if old_element["kind"] == "table":
            finding = _compare_table_content(old_element, new_element)
            if finding:
                findings.append(finding)
            else:
                formatting_finding = _compare_table_formatting(old_element, new_element)
                if formatting_finding:
                    findings.append(formatting_finding)
        if old_element["kind"] == "unknown":
            findings.append(
                _finding(
                    "minor",
                    "unsupported_shape",
                    old_element["element"],
                    f"couldn't compare (type: {old_element['shape_type']})",
                )
            )

    findings.extend(_overlap_findings(matched, stray))
    findings = _apply_report_options(findings, config)

    status = _status_from_findings(findings)
    return {"index": index, "status": status, "findings": findings}


def _apply_report_options(findings: list[dict[str, Any]], config: dict[str, Any]) -> list[dict[str, Any]]:
    normalized: list[dict[str, Any]] = []
    mode = config.get("mode", "strict")
    quiet = bool(config.get("quiet", False))

    for finding in findings:
        finding = dict(finding)
        if mode == "lenient" and finding["severity"] == "minor":
            finding["severity"] = "accepted"
        if quiet and finding["severity"] in {"minor", "accepted"}:
            continue
        normalized.append(finding)
    return normalized


def _status_from_findings(findings: list[dict[str, Any]]) -> str:
    if any(finding["severity"] == "fail" for finding in findings):
        return "fail"
    if any(finding["severity"] == "minor" for finding in findings):
        return "warning"
    return "ok"


def _compare_placement(
    old_element: dict[str, Any], new_element: dict[str, Any], slide_width: int, slide_height: int, config: dict[str, Any]
) -> dict[str, str] | None:
    if _off_slide(new_element["bounds"], slide_width, slide_height):
        return _finding("fail", "off_slide", old_element["element"], "element extends past slide bounds")

    shift_tolerance = int(config.get("shift_tolerance", 0))
    resize_tolerance = int(config.get("resize_tolerance", 0))
    old_bounds = old_element["bounds"]
    new_bounds = new_element["bounds"]
    shift_delta = max(abs(old_bounds["left"] - new_bounds["left"]), abs(old_bounds["top"] - new_bounds["top"]))
    resize_delta = max(
        abs(_width(old_bounds) - _width(new_bounds)),
        abs(_height(old_bounds) - _height(new_bounds)),
    )
    if shift_delta <= shift_tolerance and resize_delta <= resize_tolerance:
        return None
    return _finding(
        "minor",
        "placement_changed",
        old_element["element"],
        f"element moved/resized beyond tolerance (shift {shift_delta}, resize {resize_delta})",
    )


def _metric_finding(old_element: dict[str, Any], new_element: dict[str, Any]) -> dict[str, Any]:
    old_bounds = old_element["bounds"]
    new_bounds = new_element["bounds"]
    metrics: dict[str, Any] = {
        "kind": old_element["kind"],
        "shift_delta": max(abs(old_bounds["left"] - new_bounds["left"]), abs(old_bounds["top"] - new_bounds["top"])),
        "resize_delta": max(
            abs(_width(old_bounds) - _width(new_bounds)),
            abs(_height(old_bounds) - _height(new_bounds)),
        ),
        "old_bounds": old_bounds,
        "new_bounds": new_bounds,
    }
    if old_element["kind"] == "picture":
        old_image = _decoded_image(old_element["image_blob"])
        new_image = _decoded_image(new_element["image_blob"])
        metrics["old_picture_size"] = list(old_image.size)
        metrics["new_picture_size"] = list(new_image.size)
        metrics["picture_changed_percent"] = None
        if old_image.size == new_image.size:
            metrics["picture_changed_percent"] = _changed_pixel_percent(old_image, new_image)
    return {
        "severity": "ok",
        "type": "measured_metrics",
        "element": old_element["element"],
        "message": "calibration metrics",
        "metrics": metrics,
    }


def _off_slide(bounds: dict[str, int], slide_width: int, slide_height: int) -> bool:
    return bounds["left"] < 0 or bounds["top"] < 0 or bounds["right"] > slide_width or bounds["bottom"] > slide_height


def _width(bounds: dict[str, int]) -> int:
    return bounds["right"] - bounds["left"]


def _height(bounds: dict[str, int]) -> int:
    return bounds["bottom"] - bounds["top"]


def _overlap_findings(
    matched: list[tuple[dict[str, Any], dict[str, Any]]], stray: list[dict[str, Any]]
) -> list[dict[str, str]]:
    # The overlap graph is judged on the *new* slide. Matched elements carry their
    # old counterpart (used to tell whether an overlap already existed in the
    # reference); stray/inserted elements have no counterpart, so any overlap they
    # introduce is by definition new.
    nodes: list[tuple[dict[str, Any] | None, dict[str, Any]]] = [(old, new) for old, new in matched]
    nodes.extend((None, element) for element in stray)

    findings: list[dict[str, str]] = []
    for left_index, (old_left, new_left) in enumerate(nodes):
        for old_right, new_right in nodes[left_index + 1 :]:
            if not _bounds_overlap(new_left["bounds"], new_right["bounds"]):
                continue
            if new_left["kind"] == "picture" and new_right["kind"] == "picture":
                findings.append(
                    _finding(
                        "fail",
                        "picture_overlap",
                        f"{new_left['element']} / {new_right['element']}",
                        f"pictures overlap: '{new_left['name']}' and '{new_right['name']}'",
                    )
                )
            elif (
                old_left is None
                or old_right is None
                or not _bounds_overlap(old_left["bounds"], old_right["bounds"])
            ):
                findings.append(
                    _finding(
                        "fail",
                        "new_overlap",
                        f"{new_left['element']} / {new_right['element']}",
                        f"new overlap between '{new_left['name']}' and '{new_right['name']}'",
                    )
                )
    return findings


def _bounds_overlap(left: dict[str, int], right: dict[str, int]) -> bool:
    return (
        left["left"] < right["right"]
        and left["right"] > right["left"]
        and left["top"] < right["bottom"]
        and left["bottom"] > right["top"]
    )


def _match_elements(
    old_elements: list[dict[str, Any]], new_elements: list[dict[str, Any]]
) -> tuple[list[tuple[dict[str, Any], dict[str, Any]]], list[dict[str, Any]], list[dict[str, Any]]]:
    if _can_match_by_stable_id(old_elements) and _can_match_by_stable_id(new_elements):
        return _match_by_stable_id(old_elements, new_elements)
    return _match_by_kind_position_and_content(old_elements, new_elements)


def _can_match_by_stable_id(elements: list[dict[str, Any]]) -> bool:
    stable_ids = [element["stable_id"] for element in elements]
    return bool(stable_ids) and all(stable_ids) and len(stable_ids) == len(set(stable_ids))


def _match_by_stable_id(
    old_elements: list[dict[str, Any]], new_elements: list[dict[str, Any]]
) -> tuple[list[tuple[dict[str, Any], dict[str, Any]]], list[dict[str, Any]], list[dict[str, Any]]]:
    old_by_id = {element["stable_id"]: element for element in old_elements}
    new_by_id = {element["stable_id"]: element for element in new_elements}
    shared_ids = sorted(old_by_id.keys() & new_by_id.keys())
    matched = [(old_by_id[stable_id], new_by_id[stable_id]) for stable_id in shared_ids]
    missing = [old_by_id[stable_id] for stable_id in sorted(old_by_id.keys() - new_by_id.keys())]
    stray = [new_by_id[stable_id] for stable_id in sorted(new_by_id.keys() - old_by_id.keys())]
    return matched, missing, stray


def _match_by_kind_position_and_content(
    old_elements: list[dict[str, Any]], new_elements: list[dict[str, Any]]
) -> tuple[list[tuple[dict[str, Any], dict[str, Any]]], list[dict[str, Any]], list[dict[str, Any]]]:
    matched: list[tuple[dict[str, Any], dict[str, Any]]] = []
    missing: list[dict[str, Any]] = []
    stray: list[dict[str, Any]] = []
    kinds = sorted({element["kind"] for element in [*old_elements, *new_elements]})

    for kind in kinds:
        old_group = [element for element in old_elements if element["kind"] == kind]
        new_group = [element for element in new_elements if element["kind"] == kind]
        group_matches, group_missing, group_stray = _optimal_group_match(old_group, new_group)
        matched.extend(group_matches)
        missing.extend(group_missing)
        stray.extend(group_stray)

    return matched, missing, stray


def _optimal_group_match(
    old_group: list[dict[str, Any]], new_group: list[dict[str, Any]]
) -> tuple[list[tuple[dict[str, Any], dict[str, Any]]], list[dict[str, Any]], list[dict[str, Any]]]:
    if not old_group:
        return [], [], list(new_group)
    if not new_group:
        return [], list(old_group), []

    # Bounded greedy assignment within a type-class: enumerate only *eligible*
    # candidate pairs (polynomial, never factorial), sort by ascending cost, and
    # claim each old/new exactly once. A pair is eligible when the content matches
    # (a moved-but-identical element) or the regions still overlap (same slot, new
    # content). Ineligible pairs are never forced together, so a vanished element
    # replaced by an unrelated one elsewhere is reported as missing + stray rather
    # than masquerading as a content change.
    candidates: list[tuple[int, int, int]] = []
    for old_index, old_element in enumerate(old_group):
        for new_index, new_element in enumerate(new_group):
            if not _match_eligible(old_element, new_element):
                continue
            candidates.append((_match_cost(old_element, new_element), old_index, new_index))
    candidates.sort()

    matched_old: set[int] = set()
    matched_new: set[int] = set()
    matched: list[tuple[dict[str, Any], dict[str, Any]]] = []
    for _cost, old_index, new_index in candidates:
        if old_index in matched_old or new_index in matched_new:
            continue
        matched_old.add(old_index)
        matched_new.add(new_index)
        matched.append((old_group[old_index], new_group[new_index]))

    missing = [element for index, element in enumerate(old_group) if index not in matched_old]
    stray = [element for index, element in enumerate(new_group) if index not in matched_new]
    return matched, missing, stray


def _match_eligible(old_element: dict[str, Any], new_element: dict[str, Any]) -> bool:
    if old_element["content_key"] == new_element["content_key"]:
        return True
    return _bounds_overlap(old_element["bounds"], new_element["bounds"])


def _match_cost(old_element: dict[str, Any], new_element: dict[str, Any]) -> int:
    return _bounds_distance(old_element["bounds"], new_element["bounds"]) + _content_distance(old_element, new_element)


def _bounds_distance(old_bounds: dict[str, int], new_bounds: dict[str, int]) -> int:
    return sum(abs(old_bounds[key] - new_bounds[key]) for key in ("left", "top", "right", "bottom"))


def _content_distance(old_element: dict[str, Any], new_element: dict[str, Any]) -> int:
    return 0 if old_element["content_key"] == new_element["content_key"] else 10_000_000_000


def _visible_elements(shapes: Iterable[Any], config: dict[str, Any]) -> list[dict[str, Any]]:
    elements: list[dict[str, Any]] = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            elements.extend(_visible_elements(shape.shapes, config))
            continue
        element = _element(shape)
        if not _ignored(element, config):
            elements.append(element)
    return elements


def _element(shape: Any) -> dict[str, Any]:
    kind = _kind(shape)
    name = getattr(shape, "name", "") or "<unnamed>"
    stable_id = _stable_identifier(shape, name)
    shape_type = str(shape.shape_type)
    element_id = f"{kind}:{name}"
    return {
        "identity": element_id,
        "element": element_id,
        "kind": kind,
        "name": name,
        "stable_id": stable_id,
        "shape_type": shape_type,
        "content_key": _content_key(shape, kind),
        "text": shape.text_frame.text if kind == "text" else "",
        "table_cells": _table_cells(shape, kind),
        "table_formatting": _table_formatting(shape, kind),
        "text_formatting": _text_formatting(shape, kind),
        "image_blob": _image_blob(shape, kind),
        "bounds": {
            "left": int(shape.left),
            "top": int(shape.top),
            "right": int(shape.left + shape.width),
            "bottom": int(shape.top + shape.height),
        },
    }


def _stable_identifier(shape: Any, name: str) -> str | None:
    if name and name != "<unnamed>" and not _is_default_shape_name(name):
        return name
    alt_text = _alt_text(shape)
    if alt_text:
        return alt_text
    return None


def _compare_picture_content(
    old_element: dict[str, Any], new_element: dict[str, Any], config: dict[str, Any]
) -> dict[str, str] | None:
    old_image = _decoded_image(old_element["image_blob"])
    new_image = _decoded_image(new_element["image_blob"])
    # A blank candidate picture is always a regression (PRD: blank pictures FAIL),
    # whether or not the reference happened to be blank too.
    if _is_blank(new_image):
        detail = "" if _is_blank(old_image) else " where reference had content"
        return _finding("fail", "picture_blank", old_element["element"], f"picture is blank{detail}")

    if old_image.size != new_image.size:
        width_delta = abs(old_image.width - new_image.width)
        height_delta = abs(old_image.height - new_image.height)
        tolerance = int(config.get("picture_dimension_tolerance", 0))
        if not config.get("picture_normalize_resize", False) and max(width_delta, height_delta) > tolerance:
            return _finding(
                "fail",
                "picture_dimension_changed",
                old_element["element"],
                f"picture dimensions changed from {old_image.size[0]}x{old_image.size[1]} to {new_image.size[0]}x{new_image.size[1]}",
            )
        new_image = new_image.resize(old_image.size)

    changed_percent = _changed_pixel_percent(old_image, new_image)
    if changed_percent == 0:
        return None
    threshold = float(config.get("picture_pixel_threshold", 0.0))
    severity = "minor" if changed_percent <= threshold else "fail"
    return _finding(
        severity,
        "picture_pixels_changed",
        old_element["element"],
        f"picture pixels changed by {changed_percent:.2f}% (threshold {threshold:.2f}%)",
    )


def _compare_text_content(
    old_element: dict[str, Any], new_element: dict[str, Any], config: dict[str, Any]
) -> dict[str, str] | None:
    old_text = _normalize_free_text(old_element["text"], config)
    new_text = _normalize_free_text(new_element["text"], config)
    if old_text == new_text:
        return None
    return _finding("fail", "text_changed", old_element["element"], "text wording changed after volatile-token normalization")


def _compare_text_formatting(old_element: dict[str, Any], new_element: dict[str, Any]) -> dict[str, str] | None:
    if old_element["text_formatting"] == new_element["text_formatting"]:
        return None
    return _finding("minor", "text_formatting_changed", old_element["element"], "text formatting changed")


def _compare_table_content(old_element: dict[str, Any], new_element: dict[str, Any]) -> dict[str, str] | None:
    old_cells = old_element["table_cells"]
    new_cells = new_element["table_cells"]
    old_shape = _table_shape(old_cells)
    new_shape = _table_shape(new_cells)
    if old_shape != new_shape:
        return _finding(
            "fail",
            "table_structure_changed",
            old_element["element"],
            f"table structure changed from {old_shape[0]}x{old_shape[1]} to {new_shape[0]}x{new_shape[1]}",
        )
    if _has_any_data(old_cells) and not _has_any_data(new_cells):
        return _finding("fail", "table_blank", old_element["element"], "table is blank where reference had data")
    for row_index, (old_row, new_row) in enumerate(zip(old_cells, new_cells, strict=True)):
        if _has_any_data([old_row]) and not _has_any_data([new_row]):
            return _finding("fail", "table_empty_row", old_element["element"], f"table row {row_index + 1} is empty")
    for row_index, old_row in enumerate(old_cells):
        for column_index, old_text in enumerate(old_row):
            if old_text != new_cells[row_index][column_index]:
                return _finding(
                    "fail",
                    "table_cell_changed",
                    old_element["element"],
                    f"table cell R{row_index + 1}C{column_index + 1} changed",
                )
    return None


def _compare_table_formatting(old_element: dict[str, Any], new_element: dict[str, Any]) -> dict[str, str] | None:
    if old_element["table_formatting"] == new_element["table_formatting"]:
        return None
    return _finding("minor", "table_formatting_changed", old_element["element"], "table formatting changed")


def _table_shape(cells: list[list[str]]) -> tuple[int, int]:
    return len(cells), len(cells[0]) if cells else 0


def _has_any_data(cells: list[list[str]]) -> bool:
    return any(text for row in cells for text in row)


def _normalize_free_text(text: str, config: dict[str, Any]) -> str:
    patterns = config.get(
        "volatile_text_patterns",
        [
            r"\b\d{4}-\d{2}-\d{2}\b",
            r"\b\d{8}(_\d{6})?\b",
            r"\b[\w.-]+_\d{8}(_\d{6})?\.pptx\b",
        ],
    )
    normalized = text.strip()
    for pattern in patterns:
        normalized = re.sub(pattern, "", normalized)
    return re.sub(r"\s+", " ", normalized).strip()


def _decoded_image(blob: bytes) -> Image.Image:
    return Image.open(BytesIO(blob)).convert("RGBA")


def _is_blank(image: Image.Image) -> bool:
    pixels = np.asarray(image)
    if pixels[:, :, 3].max() == 0:
        return True
    return bool(np.all(pixels[:, :, :3] == 255))


def _changed_pixel_percent(old_image: Image.Image, new_image: Image.Image) -> float:
    old_pixels = np.asarray(old_image)
    new_pixels = np.asarray(new_image)
    changed = np.any(old_pixels != new_pixels, axis=2)
    return float(changed.sum() * 100 / changed.size)


def _alt_text(shape: Any) -> str | None:
    try:
        c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
    except (AttributeError, IndexError):
        return None
    return c_nv_pr.get("descr") or c_nv_pr.get("title") or None


def _is_default_shape_name(name: str) -> bool:
    return bool(re.match(r"^(Picture|TextBox|Table|Rectangle|Freeform|Group|Chart|GraphicFrame) \d+$", name))


def _content_key(shape: Any, kind: str) -> str:
    if kind == "picture":
        return getattr(getattr(shape, "image", None), "sha1", "")
    if kind == "table":
        return "\n".join("\t".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows)
    if kind == "text":
        return shape.text_frame.text.strip()
    return ""


def _table_cells(shape: Any, kind: str) -> list[list[str]]:
    if kind != "table":
        return []
    return [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]


def _table_formatting(shape: Any, kind: str) -> tuple[tuple[tuple[Any, ...], ...], ...]:
    if kind != "table":
        return ()
    rows = []
    for row in shape.table.rows:
        rows.append(tuple(_cell_formatting(cell) for cell in row.cells))
    return tuple(rows)


def _cell_formatting(cell: Any) -> tuple[Any, ...]:
    runs = tuple(_run_formatting(run) for paragraph in cell.text_frame.paragraphs for run in paragraph.runs)
    return (_fill_color(cell), runs)


def _fill_color(cell: Any) -> str | None:
    try:
        return str(cell.fill.fore_color.rgb)
    except (AttributeError, TypeError):
        return None


def _text_formatting(shape: Any, kind: str) -> tuple[tuple[Any, ...], ...]:
    if kind != "text":
        return ()
    return tuple(_run_formatting(run) for paragraph in shape.text_frame.paragraphs for run in paragraph.runs)


def _run_formatting(run: Any) -> tuple[Any, ...]:
    font = run.font
    return (font.name, font.size, font.bold, font.italic, _font_color(font))


def _font_color(font: Any) -> str | None:
    try:
        return str(font.color.rgb)
    except (AttributeError, TypeError):
        return None


def _image_blob(shape: Any, kind: str) -> bytes | None:
    if kind == "picture":
        return getattr(getattr(shape, "image", None), "blob", None)
    return None


def _kind(shape: Any) -> str:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
        return "text"
    if _looks_like_decoration(shape):
        return "decoration"
    return "unknown"


def _looks_like_decoration(shape: Any) -> bool:
    name = (getattr(shape, "name", "") or "").lower()
    return any(token in name for token in ("logo", "footer", "slide number"))


def _ignored(element: dict[str, Any], config: dict[str, Any]) -> bool:
    for rule in config.get("ignore_list", []):
        if not isinstance(rule, dict):
            continue
        if rule.get("name") and rule["name"] != element["name"]:
            continue
        if rule.get("type") and rule["type"] != element["kind"]:
            continue
        if rule.get("region") and not _inside_region(element["bounds"], rule["region"]):
            continue
        return True
    return False


def _inside_region(bounds: dict[str, int], region: dict[str, int]) -> bool:
    return (
        bounds["left"] >= region.get("left", bounds["left"])
        and bounds["top"] >= region.get("top", bounds["top"])
        and bounds["right"] <= region.get("right", bounds["right"])
        and bounds["bottom"] <= region.get("bottom", bounds["bottom"])
    )


def _finding(severity: str, finding_type: str, element: str, message: str) -> dict[str, str]:
    return {"severity": severity, "type": finding_type, "element": element, "message": message}


def _file_error(key: str, old: str | None, new: str | None, message: str) -> dict[str, Any]:
    return {
        "key": key,
        "old": old,
        "new": new,
        "status": "fail",
        "errors": [{"severity": "fail", "type": "file_error", "message": message}],
        "slides": [],
    }


def _with_summary(files: list[dict[str, Any]]) -> dict[str, Any]:
    summary = {"files": len(files), "ok": 0, "warnings": 0, "failures": 0}
    for file_result in files:
        if file_result["status"] == "fail":
            summary["failures"] += 1
        elif file_result["status"] == "warning":
            summary["warnings"] += 1
        else:
            summary["ok"] += 1
    return {"summary": summary, "files": files}
